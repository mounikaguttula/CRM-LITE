import collections 
import collections.abc
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()

# Slide 1: Title
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "2026 ITServe Synergy\nStartup Cube\nProgram Progress Update"
subtitle.text = "ITSERVE ALLIANCE • CHICAGO CHAPTER MEETING\n\nCRM Lite\nAI-Powered Metadata CRM Platform\n\nINNOVATE. PITCH. SCALE."

# Function to add a generic text slide
def add_bullet_slide(prs, title_text, bullet_points):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = title_text
    tf = body_shape.text_frame
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.text = point
        else:
            p = tf.add_paragraph()
            p.text = point

# Slide 2: The Problem
add_bullet_slide(prs, "The problem\nCRM should accelerate revenue — not become another system to manage.", [
    "RIGID SYSTEMS: Traditional CRM platforms often require schema changes, configuration work and specialized administration when a business changes.",
    "TOO MANY CLICKS: Reps spend time navigating objects, forms, reports and screens for simple questions and routine actions.",
    "AI WITHOUT CONTEXT: Generic AI assistants can struggle to understand CRM relationships, permissions and the actions a user is actually allowed to perform.",
    "Result: slower adoption, more operational overhead, and less time selling."
])

# Slide 3: Our Solution
add_bullet_slide(prs, "OUR SOLUTION: One lightweight CRM. One conversation.", [
    "CRM combines the essential sales workflow — Lead, Company, Contact, Deal, and Product — with a conversational AI layer, so any team member can capture a lead, create a deal, or check an company by simply typing what they want.",
    "Traditional CRM: CRM → Company→ New → Enter Details → Save",
    "CRM (AI Prompt): \"Create an Company for ABC Technologies.\"",
    "THE PULSE FLOW: Lead → Company + Contact → Deal → Products",
    "Single application — no separate lead-gen or CRM tools required."
])

# Slide 4: What is CRM Lite
add_bullet_slide(prs, "WHAT IS CRM LITE\nA flexible, AI-powered CRM built for modern enterprise growth", [
    "WHO IT'S FOR:",
    "- Fast-growing SMBs & Enterprises",
    "- Teams needing zero-code custom fields & objects",
    "- Organizations leveraging AI for workflow automation",
    "- Multi-tenant SaaS companies",
    "TEN CORE CAPABILITIES:",
    "- Anthropic Claude AI, Multi-Tenant Security, Dynamic Metadata, Role Permissions (RBAC), Universal Data Engine, Field-Level Security, Lead QR Scanner, Flow Automations, Campaign Management, Self-Healing Validation",
    "Instant Customization: Define custom objects & fields on-the-fly without code changes.",
    "Claude AI Assistant: Native MCP integration — Claude manages CRM records via conversation."
])

# Slide 5: Executive Summary
add_bullet_slide(prs, "EXECUTIVE SUMMARY\nOne platform. Every deal-facing workflow.", [
    "01 Metadata-Driven Core: Admins create objects, fields, and validation rules from the UI — no engineering backlog required.",
    "02 AI Built In, Not Bolted On: A native Claude-powered chat panel lets any user query, summarize, and act on CRM data conversationally.",
    "03 Full Sales Lifecycle: Leads, Companies, Contacts, Deals, Line Items, Campaigns, and capture forms in one connected system.",
    "Positioning: everything a growing sales team needs from Salesforce-class CRM — reimagined as a fast, modern React application."
])

# Slide 6: Key Differentiator - Natural Language
add_bullet_slide(prs, "KEY DIFFERENTIATOR\nCRM through natural language", [
    "Instead of navigating screens, users simply say what they want. CRM Lite understands, executes, and confirms the action — turning CRM into a conversation.",
    "MCP integration: Users can access and perform CRM operations directly from Claude Desktop.",
    "Examples:",
    "- “Create an Company for ABC Technologies.” → Creates Company",
    "- “Create a Contact John Smith for ABC Technologies.” → Creates & links Contact",
    "- “Create a Deal for ABC Technologies worth $100,000.” → Creates Deal",
    "- “What is the Deal amount for ABC Technologies?” → Retrieves Deal info",
    "- “Show me the products on the ABC Technologies Deal.” → Retrieves Products",
    "- “Set the Deal amount to $50,000.” → Updates Deal"
])

# Slide 7: Core Platform
add_bullet_slide(prs, "CORE PLATFORM\nOne connected system, from first click to closed deal", [
    "Core Sales Objects: Lead · Company · Contact · Deal · Line Items",
    "- Full CRUD across every object",
    "- Leads convert into Companies, Contacts & Deals",
    "- Line Items capture products and pricing per deal",
    "Demand Capture: Lead Capture Forms · Campaigns",
    "- Web-to-lead forms with built-in CAPTCHA protection",
    "- Campaigns tracked as first-class records",
    "- Leads linked back to the campaign that sourced them"
])

# Slide 8: Key Differentiator - AI Teammate
add_bullet_slide(prs, "KEY DIFFERENTIATOR\nAn AI teammate inside every workspace", [
    "A chat panel connected to Claude sits alongside every record, so reps and managers can ask questions and get action in natural language — instead of building reports or hunting through screens.",
    "- Summarize a deal or account in seconds",
    "- Answer questions about pipeline and history",
    "- Draft follow-ups grounded in real CRM data",
    "- Available from any object, any workspace"
])

# Slide 9: Lead Capture
add_bullet_slide(prs, "LEAD CAPTURE\nBuilt-in forms, published everywhere", [
    "Create Form → Generate Link → Publish Across Channels → Prospect Submits → Lead Captured → Qualified & Converted",
    "Channels supported at launch:",
    "- LinkedIn, Company Website, Social Media",
    "- Email Campaigns, Digital Ads, Online Communities",
    "WHY IT MATTERS:",
    "Marketing and Sales share one system. Leads generated from any channel flow directly into the pipeline — no manual re-entry, no disconnected lead-gen tools."
])

# Slide 10 & 11: Claude + MCP
add_bullet_slide(prs, "Claude + Model Context Protocol", [
    "CRM actions are available through a conversational interface.",
    "USER: Types a request in natural language. Examples: create, retrieve, update and manage CRM records.",
    "MCP TOOL LAYER: Claude uses connected MCP tools to interact with the CRM's structured actions and data.",
    "CRM LITE: Executes the permitted operation and returns the result to the user.",
    "CRM from Claude Desktop: One CRM data layer. Multiple AI entry points."
])

# Slide 12: Differentiators
add_bullet_slide(prs, "Seven differentiators that set us apart", [
    "1 AI-Powered CRM — natural-language actions, not just chat",
    "2 Lightweight — essential capability, zero bloat",
    "3 Built-in lead capture — no separate form tool needed",
    "4 Multi-channel forms — one link, every channel",
    "5 Complete Lead → Deal journey in one app",
    "6 Conversational data access — ask, don't click",
    "7 Multi-user SaaS — collaborative from day one"
])

# Slide 13: Technology & Architecture
add_bullet_slide(prs, "TECHNOLOGY & ARCHITECTURE\nA modern, maintainable stack", [
    "React 19 Frontend: Dynamic, metadata-driven UI rendering + chat panel",
    "↔ Express.js Backend: Self-healing validation rules, multi-tenant routing",
    "↔ Supabase PostgreSQL + AI: Universal JSONB engine + Claude AI via MCP",
    "Why it's differentiated technically:",
    "- Universal JSONB engine means new objects & fields ship instantly — zero database migrations",
    "- Claude connects via MCP with OAuth 2.1, so AI access always respects RBAC and field-level security",
    "- Self-healing validation rules and multi-tenant isolation are built into the data layer, not bolted on"
])

# Slide 14: Who we are building for
add_bullet_slide(prs, "Who we are building for\nStart with teams that need CRM speed without enterprise CRM complexity.", [
    "SMBs & GROWING TEAMS: Sales teams that need the essentials quickly and want less administrative overhead.",
    "AI-FIRST ORGANIZATIONS: Organizations adopting AI for day-to-day business operations and workflow execution.",
    "MULTI-TENANT BUSINESSES: SaaS companies and consultancies that need flexible CRM structures across organizations.",
    "Positioning: lightweight enough to adopt quickly, flexible enough to adapt, intelligent enough to operate conversationally."
])

# Slide 15: Roadmap
add_bullet_slide(prs, "ROADMAP\nWhat's next", [
    "Near-Term: Expand AI actions — let the chat panel create and update records, not just answer questions.",
    "Mid-Term: Deeper analytics and reporting layer across all metadata-driven objects.",
    "Ongoing: Broaden integrations (email, calendar) and mobile-friendly workspace access."
])

# Slide 16: Competitive Advantage
add_bullet_slide(prs, "THE COMPETITIVE ADVANTAGE\nWhy CRM Lite stands out", [
    "Bespoke Flexibility: Delivers custom software adaptability without expensive development bills or code rewrites. Define new objects, fields, and rules instantly.",
    "Native Claude AI Assistant: Built from the ground up for Anthropic Claude AI integration via Model Context Protocol (MCP). Natural language CRM management and data operations.",
    "Enterprise Governance: Built-in multi-tenancy, granular RBAC, Field-Level Security, and self-healing validation rules — production-ready out of the box.",
    "\"Bridging the gap between rigid SaaS tools and custom software — empowering organizations with total schema flexibility and autonomous Claude AI.\""
])

# Slide 17: Vision
add_bullet_slide(prs, "OUR VISION\nCRM should feel like a conversation, not a chore.", [
    "CRM's long-term vision is an AI-first CRM where users never have to remember navigation — they simply say what they want to accomplish, and the platform gets it done.",
    "[founder@pulsecrm.com] | [www.pulsecrm.com] | [LinkedIn / Company Handle]"
])


prs.save('c:\\Users\\Lenovo\\LiteWeight-CRM\\ITServe_Synergy_Startup_Cube_CRM_Lite.pptx')
print("Presentation generated!")
